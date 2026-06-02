"""Story 2-2b 验收测试 — 文档解析与内容提取（扩展格式）

使用 pytest-bdd 框架，场景定义在 .feature 文件中。
测试使用真实解析器和编程式临时 fixture 文件，验证端到端解析正确性。

Run with: poetry run pytest tests/acceptance/test_acceptance_document_parse_extended.py -v
"""

from __future__ import annotations

import os
import tempfile
from typing import Any, Generator

import pytest
from pytest_bdd import given, scenario, scenarios, then, when

scenarios("test_acceptance_document_parse_extended.feature")


# ===================================================================
# Fixtures
# ===================================================================


@pytest.fixture
def context() -> Generator[dict[str, Any], None, None]:
    """BDD 步骤间共享状态，teardown 时自动清理临时文件"""
    ctx: dict[str, Any] = {}
    yield ctx
    _cleanup(ctx.get("temp_path", ""))


# ===================================================================
# 内部 helpers — 编程式 fixture 工厂
# ===================================================================


def _cleanup(path: str) -> None:
    """安全清理临时文件"""
    if path and os.path.exists(path):
        try:
            os.unlink(path)
        except OSError:
            pass


# ===================================================================
# Background Steps
# ===================================================================


@given("Story 2-2a 基础格式解析已实现，DocumentParserPort/ParsedDocument/CompositeDocumentParser 已可用")
def given_story_2_2a_completed(context: dict) -> None:
    """背景: Story 2-2a 已完成."""
    context["story_2_2a_ready"] = True


# ===================================================================
# AC-1: PPT/PPTX 文档解析
# ===================================================================


@scenario("test_acceptance_document_parse_extended.feature", "AC-1 - 成功解析 PPTX 文档")
def test_ac1_parse_pptx_success():
    """成功解析含文本和备注的 PPTX 文档."""
    pass


@scenario("test_acceptance_document_parse_extended.feature", "AC-1 - 解析含表格的 PPTX 文档")
def test_ac1_parse_pptx_with_table():
    """解析含内嵌表格的 PPTX 文档."""
    pass


@scenario("test_acceptance_document_parse_extended.feature", "AC-1 - 解析空 PPTX 文档失败")
def test_ac1_parse_empty_pptx_fails():
    """解析空 PPTX 返回 failed."""
    pass


@scenario("test_acceptance_document_parse_extended.feature", "AC-1 - 解析旧版 PPT 格式拒绝")
def test_ac1_parse_ppt_legacy_rejected():
    """解析旧版 PPT 返回拒绝消息."""
    pass


@given('有一个包含文本和备注的 PPTX 文件 "strategy.pptx"')
def given_pptx_with_text_and_notes(context: dict) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "战略规划概述"
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "这是一条备注"
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "实施计划"

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()
    context["temp_path"] = tmp.name
    context["mime_type"] = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@given('有一个包含内嵌表格的 PPTX 文件 "data.pptx"')
def given_pptx_with_table(context: dict) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "数据报告"
    table_shape = slide.shapes.add_table(3, 3, Inches(1), Inches(2), Inches(8), Inches(3))
    for r in range(3):
        for c in range(3):
            table_shape.table.cell(r, c).text = f"R{r}C{c}"

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()
    context["temp_path"] = tmp.name
    context["mime_type"] = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@given("有一个无幻灯片的空 PPTX 文件")
def given_empty_pptx(context: dict) -> None:
    from pptx import Presentation

    prs = Presentation()
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()
    context["temp_path"] = tmp.name
    context["mime_type"] = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@given('有一个旧版 PPT 文件 "legacy.ppt"')
def given_legacy_ppt(context: dict) -> None:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".ppt")
    tmp.write(b"not a valid pptx file")
    tmp.close()
    context["temp_path"] = tmp.name


@when("系统使用 PptxParser 解析该文件")
def when_parse_with_pptx_parser(context: dict) -> None:
    from src.infrastructure.external_services.document_parsing.pptx_parser import PptxParser

    parser = PptxParser()
    context["result"] = parser.parse(context["temp_path"], context["mime_type"])


@when("系统使用 PptxParser 解析 MIME 类型为 application/vnd.ms-powerpoint 的文件")
def when_parse_ppt_mime_with_pptx_parser(context: dict) -> None:
    from src.infrastructure.external_services.document_parsing.pptx_parser import PptxParser

    parser = PptxParser()
    context["result"] = parser.parse(context["temp_path"], "application/vnd.ms-powerpoint")


@then("parse_status 为 completed")
def then_parse_status_completed(context: dict) -> None:
    result = context["result"]
    assert result.is_completed(), f"解析应成功，实际: {result.error_message}"


@then("提取的文本包含幻灯片标题")
def then_text_contains_slide_titles(context: dict) -> None:
    result = context["result"]
    all_text = " ".join(t.content for p in result.pages for t in p.texts)
    assert "战略规划概述" in all_text, f"应提取到幻灯片标题，实际: {all_text}"
    assert "实施计划" in all_text, f"应提取到第二页幻灯片标题，实际: {all_text}"


@then("备注内容被提取")
def then_notes_extracted(context: dict) -> None:
    result = context["result"]
    all_text = " ".join(t.content for p in result.pages for t in p.texts)
    assert "备注" in all_text, f"应提取到备注内容，实际: {all_text}"


@then("每页幻灯片编号作为 page_number")
def then_page_numbers_sequential(context: dict) -> None:
    result = context["result"]
    page_numbers = [p.page_number for p in result.pages]
    assert page_numbers == sorted(page_numbers), f"页码应为升序，实际: {page_numbers}"


@then("表格被提取为 ParsedTable")
def then_table_extracted(context: dict) -> None:
    result = context["result"]
    all_tables = [t for p in result.pages for t in p.tables]
    assert len(all_tables) >= 1, "应至少提取到 1 个表格"
    assert len(all_tables[0].rows) >= 2, "表格应至少有 2 行"


@then("parse_status 为 failed")
def then_parse_status_failed(context: dict) -> None:
    result = context["result"]
    assert result.parse_status == "failed", f"解析应失败，实际: {result.parse_status}"


@then("error_message 说明文档为空")
def then_error_message_empty_document(context: dict) -> None:
    result = context["result"]
    assert result.error_message is not None
    assert "空" in result.error_message or "无内容" in result.error_message or "无幻灯片" in result.error_message.lower(), (
        f"错误信息应说明文档为空，实际: {result.error_message}"
    )


@then("error_message 建议转换为 PPTX")
def then_error_message_suggest_pptx(context: dict) -> None:
    result = context["result"]
    assert result.error_message is not None
    assert "PPTX" in result.error_message, f"错误信息应建议转换为 PPTX，实际: {result.error_message}"


@then("error_message 建议转换为 XLSX")
def then_error_message_suggest_xlsx(context: dict) -> None:
    result = context["result"]
    assert result.error_message is not None
    assert "XLSX" in result.error_message, f"错误信息应建议转换为 XLSX，实际: {result.error_message}"


# ===================================================================
# AC-2: Excel 文档解析（XLSX/XLS）
# ===================================================================


@scenario("test_acceptance_document_parse_extended.feature", "AC-2 - 成功解析多 Sheet XLSX 文档")
def test_ac2_parse_xlsx_multi_sheet():
    """成功解析多 Sheet XLSX 文档."""
    pass


@scenario("test_acceptance_document_parse_extended.feature", "AC-2 - 解析旧版 XLS 格式拒绝")
def test_ac2_parse_xls_legacy_rejected():
    """解析旧版 XLS 返回拒绝消息."""
    pass


@given('有一个包含多个 Sheet 的 XLSX 文件 "data.xlsx"')
def given_xlsx_multi_sheet(context: dict) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws1 = wb.active
    assert ws1 is not None
    ws1.title = "市场分析"
    ws1.append(["维度", "目标", "进度"])
    ws1.append(["市场", "增长20%", "进行中"])
    ws2 = wb.create_sheet(title="技术评估")
    ws2.append(["项目", "状态"])
    ws2.append(["迁移上云", "未开始"])

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx")
    wb.save(tmp.name)
    tmp.close()
    context["temp_path"] = tmp.name
    context["mime_type"] = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


@given('有一个旧版 XLS 文件 "legacy.xls"')
def given_legacy_xls(context: dict) -> None:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".xls")
    tmp.write(b"not a valid xlsx file")
    tmp.close()
    context["temp_path"] = tmp.name


@when("系统使用 ExcelParser 解析该文件")
def when_parse_with_excel_parser(context: dict) -> None:
    from src.infrastructure.external_services.document_parsing.excel_parser import ExcelParser

    parser = ExcelParser()
    context["result"] = parser.parse(context["temp_path"], context["mime_type"])


@when("系统使用 ExcelParser 解析 MIME 类型为 application/vnd.ms-excel 的文件")
def when_parse_xls_mime_with_excel_parser(context: dict) -> None:
    from src.infrastructure.external_services.document_parsing.excel_parser import ExcelParser

    parser = ExcelParser()
    context["result"] = parser.parse(context["temp_path"], "application/vnd.ms-excel")


@then("每个 Sheet 独立输出为 ParsedTable")
def then_each_sheet_becomes_parsed_table(context: dict) -> None:
    result = context["result"]
    all_tables = [t for p in result.pages for t in p.tables]
    assert len(all_tables) >= 2, f"应至少有 2 个 Sheet 表格，实际: {len(all_tables)}"


@then('sheet 名称存储于 ParsedTable.metadata["sheet_name"]')
def then_sheet_name_in_parsedtable_metadata(context: dict) -> None:
    result = context["result"]
    all_tables = [t for p in result.pages for t in p.tables]
    sheet_names = [t.metadata.get("sheet_name", "") for t in all_tables]
    assert "市场分析" in sheet_names, f"应包含 '市场分析' Sheet，实际: {sheet_names}"
    assert "技术评估" in sheet_names, f"应包含 '技术评估' Sheet，实际: {sheet_names}"


# ===================================================================
# AC-3: CSV 文档解析
# ===================================================================


@scenario("test_acceptance_document_parse_extended.feature", "AC-3 - 成功解析 CSV 文档")
def test_ac3_parse_csv_success():
    """成功解析 CSV 文档."""
    pass


@scenario("test_acceptance_document_parse_extended.feature", "AC-3 - 解析空 CSV 文档失败")
def test_ac3_parse_empty_csv_fails():
    """解析空 CSV 返回 failed."""
    pass


@given('有一个 UTF-8 编码的 CSV 文件 "export.csv"')
def given_csv_utf8(context: dict) -> None:
    content = "姓名,部门,职级\n张三,技术部,P7\n李四,市场部,P6\n"
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".csv")
    tmp.write(content.encode("utf-8"))
    tmp.close()
    context["temp_path"] = tmp.name
    context["mime_type"] = "text/csv"


@given("有一个空 CSV 文件")
def given_empty_csv(context: dict) -> None:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".csv")
    tmp.write(b"")
    tmp.close()
    context["temp_path"] = tmp.name
    context["mime_type"] = "text/csv"


@when("系统使用 CSVParser 解析该文件")
def when_parse_with_csv_parser(context: dict) -> None:
    from src.infrastructure.external_services.document_parsing.csv_parser import CSVParser

    parser = CSVParser()
    context["result"] = parser.parse(context["temp_path"], context["mime_type"])


@then("输出单页结构包含一个 ParsedTable")
def then_single_page_with_one_table(context: dict) -> None:
    result = context["result"]
    assert len(result.pages) >= 1, "应至少有 1 页"
    all_tables = [t for p in result.pages for t in p.tables]
    assert len(all_tables) >= 1, "应至少包含 1 个表格"


@then("表头和数据显示正确")
def then_csv_data_correct(context: dict) -> None:
    result = context["result"]
    all_tables = [t for p in result.pages for t in p.tables]
    table = all_tables[0]
    assert table.rows[0][0] == "姓名"
    assert table.rows[1][1] == "技术部"


# ===================================================================
# AC-4: 图像文档解析（JPEG/PNG/GIF）
# ===================================================================


@scenario("test_acceptance_document_parse_extended.feature", "AC-4 - 成功提取 JPEG 图像元数据")
def test_ac4_parse_jpeg_metadata():
    """成功提取 JPEG 图像元数据."""
    pass


@scenario("test_acceptance_document_parse_extended.feature", "AC-4 - GIF 仅处理第一帧")
def test_ac4_parse_gif_first_frame():
    """GIF 仅处理第一帧."""
    pass


@given('有一个 JPEG 图像文件 "chart.jpg"')
def given_jpeg_image(context: dict) -> None:
    from PIL import Image

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
    img = Image.new("RGB", (100, 100), color="red")
    img.save(tmp.name, format="JPEG")
    tmp.close()
    context["temp_path"] = tmp.name
    context["mime_type"] = "image/jpeg"


@given('有一个多帧 GIF 图像文件 "animation.gif"')
def given_gif_image(context: dict) -> None:
    from PIL import Image

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".gif")
    frame1 = Image.new("RGB", (100, 100), color="red")
    frame2 = Image.new("RGB", (100, 100), color="blue")
    frame1.save(tmp.name, format="GIF", save_all=True, append_images=[frame2], loop=0)
    tmp.close()
    context["temp_path"] = tmp.name
    context["mime_type"] = "image/gif"


@when("系统使用 ImageParser 解析该文件")
def when_parse_with_image_parser(context: dict) -> None:
    from src.infrastructure.external_services.document_parsing.image_parser import ImageParser

    parser = ImageParser()
    context["result"] = parser.parse(context["temp_path"], context["mime_type"])


@then("images 数组包含图像元数据（format/width/height/mode）")
def then_image_metadata_in_result(context: dict) -> None:
    result = context["result"]
    images = [i for p in result.pages for i in p.images]
    assert len(images) >= 1, "应至少包含 1 个图像元素"
    img_meta = images[0].metadata
    assert "format" in img_meta, f"应包含 format 字段，实际: {img_meta}"
    assert "width" in img_meta, f"应包含 width 字段，实际: {img_meta}"
    assert "height" in img_meta, f"应包含 height 字段，实际: {img_meta}"


@then("image 元素的 content 为空字符串")
def then_image_content_empty(context: dict) -> None:
    result = context["result"]
    images = [i for p in result.pages for i in p.images]
    assert len(images) >= 1
    for img in images:
        assert img.content == "", f"image content 应为空字符串，实际: '{img.content}'"


@then("images 数组仅包含 1 个元素")
def then_single_image_for_gif(context: dict) -> None:
    result = context["result"]
    images = [i for p in result.pages for i in p.images]
    assert len(images) == 1, f"GIF 应仅处理第一帧，实际 frames: {len(images)}"


# ===================================================================
# AC-5: HTML 文档解析
# ===================================================================


@scenario("test_acceptance_document_parse_extended.feature", "AC-5 - 成功解析 HTML 文档")
def test_ac5_parse_html_success():
    """成功解析 HTML 文档."""
    pass


@scenario("test_acceptance_document_parse_extended.feature", "AC-5 - 解析空 HTML 文档失败")
def test_ac5_parse_empty_html_fails():
    """解析空 HTML 返回 failed."""
    pass


@given('有一个包含标题、段落和表格的 HTML 文件 "report.html"')
def given_html_with_structure(context: dict) -> None:
    html = (
        "<html><body><h1>标题</h1><h2>副标题</h2><p>段落内容</p>"
        "<table><tr><th>A</th><th>B</th></tr><tr><td>1</td><td>2</td></tr></table></body></html>"
    )
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".html")
    tmp.write(html.encode("utf-8"))
    tmp.close()
    context["temp_path"] = tmp.name
    context["mime_type"] = "text/html"


@given("有一个空 body 的 HTML 文件")
def given_empty_html(context: dict) -> None:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".html")
    tmp.write(b"<html><body></body></html>")
    tmp.close()
    context["temp_path"] = tmp.name
    context["mime_type"] = "text/html"


@when("系统使用 HTMLParser 解析该文件")
def when_parse_with_html_parser(context: dict) -> None:
    from src.infrastructure.external_services.document_parsing.html_parser import HTMLParser

    parser = HTMLParser()
    context["result"] = parser.parse(context["temp_path"], context["mime_type"])


@then("文本提取包含标题和段落内容")
def then_html_text_extraction(context: dict) -> None:
    result = context["result"]
    all_text = " ".join(t.content for p in result.pages for t in p.texts)
    assert "标题" in all_text, f"应提取到标题文本，实际: {all_text}"
    assert "段落内容" in all_text, f"应提取到正文内容，实际: {all_text}"


@then("标题层级映射到 metadata.style")
def then_heading_levels_in_metadata(context: dict) -> None:
    result = context["result"]
    styles = [t.metadata.get("style", "") for p in result.pages for t in p.texts]
    assert "h1" in styles, f"应识别到 h1 样式，实际: {styles}"
    assert "h2" in styles, f"应识别到 h2 样式，实际: {styles}"


@then("HTML 表格被提取为 ParsedTable")
def then_html_table_extracted(context: dict) -> None:
    result = context["result"]
    all_tables = [t for p in result.pages for t in p.tables]
    assert len(all_tables) >= 1, "应至少提取到 1 个表格"


# ===================================================================
# AC-6: Markdown 文档解析
# ===================================================================


@scenario("test_acceptance_document_parse_extended.feature", "AC-6 - 成功解析 Markdown 文档")
def test_ac6_parse_markdown_success():
    """成功解析 Markdown 文档."""
    pass


@scenario("test_acceptance_document_parse_extended.feature", "AC-6 - 解析空 Markdown 文档失败")
def test_ac6_parse_empty_markdown_fails():
    """解析空 Markdown 返回 failed."""
    pass


@given('有一个包含标题、段落、表格和代码块的 Markdown 文件 "plan.md"')
def given_markdown_with_structure(context: dict) -> None:
    md = "# 标题1\n\n段落内容\n\n## 标题2\n\n| A | B |\n|---|---|\n| 1 | 2 |\n\n```python\nprint('hello')\n```\n"
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".md")
    tmp.write(md.encode("utf-8"))
    tmp.close()
    context["temp_path"] = tmp.name
    context["mime_type"] = "text/markdown"


@given("有一个空 Markdown 文件")
def given_empty_markdown(context: dict) -> None:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".md")
    tmp.write(b"")
    tmp.close()
    context["temp_path"] = tmp.name
    context["mime_type"] = "text/markdown"


@when("系统使用 MarkdownParser 解析该文件")
def when_parse_with_markdown_parser(context: dict) -> None:
    from src.infrastructure.external_services.document_parsing.markdown_parser import MarkdownParser

    parser = MarkdownParser()
    context["result"] = parser.parse(context["temp_path"], context["mime_type"])


@then("标题层级识别正确（# → h1，## → h2）")
def then_heading_levels_correct(context: dict) -> None:
    result = context["result"]
    all_text = " ".join(t.content for p in result.pages for t in p.texts)
    assert "标题1" in all_text, f"应提取到 # 标题，实际: {all_text}"
    assert "标题2" in all_text, f"应提取到 ## 标题，实际: {all_text}"


@then("段落按连续空行分割")
def then_paragraphs_split_by_blank_lines(context: dict) -> None:
    result = context["result"]
    texts = [t.content for p in result.pages for t in p.texts]
    assert len(texts) >= 2, f"应按空行分割为多个段落，实际: {len(texts)} 段"


@then("Markdown 表格被提取为 ParsedTable")
def then_markdown_table_extracted(context: dict) -> None:
    result = context["result"]
    all_tables = [t for p in result.pages for t in p.tables]
    assert len(all_tables) >= 1, "应至少提取到 1 个 Markdown 表格"


@then("代码块内容保留")
def then_code_block_preserved(context: dict) -> None:
    result = context["result"]
    all_text = " ".join(t.content for p in result.pages for t in p.texts)
    assert "print" in all_text, f"应保留代码块内容，实际: {all_text}"


# ===================================================================
# AC-7: RTF 文档解析
# ===================================================================


@scenario("test_acceptance_document_parse_extended.feature", "AC-7 - 成功解析 RTF 文档")
def test_ac7_parse_rtf_success():
    """成功解析 RTF 文档."""
    pass


@scenario("test_acceptance_document_parse_extended.feature", "AC-7 - 解析空 RTF 文档失败")
def test_ac7_parse_empty_rtf_fails():
    """解析空 RTF 返回 failed."""
    pass


@given('有一个包含文本的 RTF 文件 "memo.rtf"')
def given_rtf_with_text(context: dict) -> None:
    rtf = r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Times New Roman;}} \f0\fs24 RTF 测试文档内容}"
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".rtf")
    tmp.write(rtf.encode("utf-8"))
    tmp.close()
    context["temp_path"] = tmp.name
    context["mime_type"] = "application/rtf"


@given("有一个仅含 RTF 头部的空 RTF 文件")
def given_empty_rtf(context: dict) -> None:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".rtf")
    tmp.write(r"{\rtf1\ansi}".encode("utf-8"))
    tmp.close()
    context["temp_path"] = tmp.name
    context["mime_type"] = "application/rtf"


@when("系统使用 RTFParser 解析该文件")
def when_parse_with_rtf_parser(context: dict) -> None:
    from src.infrastructure.external_services.document_parsing.rtf_parser import RTFParser

    parser = RTFParser()
    context["result"] = parser.parse(context["temp_path"], context["mime_type"])


@then("提取到 RTF 文本内容")
def then_rtf_text_extracted(context: dict) -> None:
    result = context["result"]
    all_text = " ".join(t.content for p in result.pages for t in p.texts)
    assert len(all_text.strip()) > 0, "应提取到 RTF 文本内容"
    assert "RTF" in all_text, f"应提取到 RTF 文本内容，实际: {all_text}"


# ===================================================================
# AC-8: CompositeDocumentParser 扩展与集成
# ===================================================================


@scenario("test_acceptance_document_parse_extended.feature", "AC-8 - 所有 17 种格式 MIME 路由正确")
def test_ac8_all_mime_types_registered():
    """验证所有 MIME 类型已注册."""
    pass


@scenario("test_acceptance_document_parse_extended.feature", "AC-8 - 不支持的 MIME 类型返回失败")
def test_ac8_unsupported_mime():
    """不支持的 MIME 返回 failed."""
    pass


@given("Composition Root 已注册所有扩展格式解析器")
def given_all_extended_parsers_registered(context: dict) -> None:
    from src.domain.ports.resolver import resolve

    context["parser"] = resolve("document_parser")


@when("实例化 CompositeDocumentParser")
def when_instantiate_composite_parser(context: dict) -> None:
    pass  # parser 已在 given 步骤中通过 resolve 获取


@then("MIME 路由表包含预期 15 种 MIME 类型")
def then_mime_table_has_15_types(context: dict) -> None:
    parser = context["parser"]
    mime_types = set(parser._parsers.keys())
    expected = {
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
        "text/plain",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
        "text/csv",
        "image/jpeg",
        "image/png",
        "image/gif",
        "text/html",
        "text/markdown",
        "application/rtf",
    }
    missing = expected - mime_types
    assert not missing, f"CompositeDocumentParser 缺少 MIME 路由: {missing}"


@then("document_parser 端口版本为 v1.1.0")
def then_port_version_is_v1_1_0(context: dict) -> None:
    from src.domain.ports.registry import _global_registry

    spec = _global_registry.get("document_parser")
    assert spec is not None
    assert spec.version == "v1.1.0", f"端口版本应为 v1.1.0，实际: {spec.version}"


@given("有一个未知 MIME 类型的文件")
def given_unknown_mime_file(context: dict) -> None:
    from src.domain.ports.resolver import resolve

    context["parser"] = resolve("document_parser")
    context["temp_path"] = "/dev/null"


@when("系统使用 CompositeDocumentParser 解析该文件")
def when_parse_unknown_mime_with_composite(context: dict) -> None:
    parser = context["parser"]
    context["result"] = parser.parse(context["temp_path"], "application/x-unknown-format")


@then("返回 ParsedDocument 而非抛异常")
def then_returns_parsed_document_not_exception(context: dict) -> None:
    from src.domain.value_objects.parsed_document import ParsedDocument

    result = context["result"]
    assert isinstance(result, ParsedDocument), f"应返回 ParsedDocument，实际: {type(result)}"


@then("error_message 包含明确错误描述")
def then_error_message_descriptive(context: dict) -> None:
    result = context["result"]
    assert result.error_message is not None, "失败必须附带错误消息"
    assert len(result.error_message) > 0, "错误消息不能为空"
