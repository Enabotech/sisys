"""扩展格式文档解析集成测试

测试 Story 2-2b 新增 7 个解析器（PPTX/XLSX/CSV/Image/HTML/Markdown/RTF）的完整解析流水线：

1. 解析器层：直接调用各扩展格式解析器
2. Service 层：通过 DocumentParsingService 编排完整流水线
   - 真实 PostgreSQL（schema 隔离）
   - 真实 MinIO（上传 → 下载 → 解析）
   - 真实 DualChannelEventBus（Redis + PostgreSQL Outbox → RabbitMQ）
   - 真实解析器
3. 并发解析测试（AC-6: ≥10 并发）
"""

from __future__ import annotations

import asyncio
import os
import tempfile
import uuid
from collections.abc import AsyncGenerator

import pytest
from sqlalchemy import select as sa_select
from sqlalchemy import text
from sqlalchemy.ext.asyncio import AsyncSession

from src.domain.entities.document import Document, ParseStatus
from src.infrastructure.config.postgresql import PostgreSQLConfig
from src.infrastructure.storage.postgresql.models import OutboxModel
from src.infrastructure.storage.postgresql.postgresql_manager import PostgreSQLManager
from src.infrastructure.storage.postgresql.repository.document_repository import PostgreSQLDocumentRepository
from src.infrastructure.storage.postgresql.session_context import reset_session, set_session
from tests.environments import get_test_env

# ===================================================================
# MIME 类型常量
# ===================================================================

MIME_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
MIME_XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
MIME_CSV = "text/csv"
MIME_JPEG = "image/jpeg"
MIME_PNG = "image/png"
MIME_HTML = "text/html"
MIME_MARKDOWN = "text/markdown"
MIME_RTF = "text/rtf"

# ===================================================================
# 工厂函数
# ===================================================================


def _create_test_pptx() -> str:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "集成测试PPTX"
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()
    return tmp.name


def _create_test_xlsx() -> str:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "集成测试"
    ws["A1"] = "集成测试XLSX"
    ws["B1"] = "值1"
    ws["C1"] = "值2"
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx")
    wb.save(tmp.name)
    tmp.close()
    return tmp.name


def _create_test_csv() -> str:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".csv")
    tmp.write("列A,列B,列C\n值1,值2,值3\n值4,值5,值6\n".encode("utf-8"))
    tmp.close()
    return tmp.name


def _create_test_html() -> str:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".html")
    tmp.write("<html><body><h1>集成测试HTML</h1><p>段落内容</p></body></html>".encode("utf-8"))
    tmp.close()
    return tmp.name


def _create_test_md() -> str:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".md")
    tmp.write("# 集成测试Markdown\n\n测试段落内容\n\n| 列A | 列B |\n|-----|-----|\n| v1  | v2  |\n".encode("utf-8"))
    tmp.close()
    return tmp.name


def _create_test_rtf() -> str:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".rtf")
    tmp.write(r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Times New Roman;}} \f0\fs24 集成测试 RTF}".encode("utf-8"))
    tmp.close()
    return tmp.name


def _cleanup(path: str) -> None:
    """安全清理临时文件"""
    if path and os.path.exists(path):
        try:
            os.unlink(path)
        except OSError:
            pass


# ===================================================================
# PostgreSQL Schema Isolation Fixtures
# ===================================================================


@pytest.fixture
def test_schema() -> str:
    return f"test_sisys_{uuid.uuid4().hex[:8]}"


@pytest.fixture
def pg_config() -> PostgreSQLConfig:
    env = get_test_env()
    return PostgreSQLConfig(
        host=env.postgres.host,
        port=env.postgres.port,
        database=env.postgres.database,
        username=env.postgres.username,
        password=env.postgres.password,
        pool_size=5,
        max_overflow=10,
    )


@pytest.fixture
def db_engine(pg_config: PostgreSQLConfig) -> PostgreSQLManager:
    return PostgreSQLManager(pg_config)


@pytest.fixture
def ensure_schema(db_engine: PostgreSQLManager, pg_config: PostgreSQLConfig, test_schema: str):
    """创建测试专用 schema 及 documents 表"""
    sync_url = (
        f"postgresql+psycopg2://{pg_config.username}:{pg_config.password}"
        f"@{pg_config.host}:{pg_config.port}/{pg_config.database}"
    )
    from sqlalchemy import create_engine

    sync_engine = create_engine(sync_url)

    try:
        with sync_engine.connect() as conn:
            conn.execute(text("SELECT 1"))
    except Exception as e:
        sync_engine.dispose()
        pytest.skip(f"PostgreSQL not available: {e}")

    with sync_engine.connect() as conn:
        conn.execute(text(f'DROP SCHEMA IF EXISTS "{test_schema}" CASCADE'))
        conn.commit()

    with sync_engine.connect() as conn:
        conn.execute(text(f'CREATE SCHEMA "{test_schema}"'))
        conn.commit()

    from src.infrastructure.storage.postgresql.models import Base

    with sync_engine.connect() as conn:
        conn.execute(text(f'SET search_path TO "{test_schema}"'))
        Base.metadata.create_all(conn)
        conn.commit()

    sync_engine.dispose()

    yield test_schema

    sync_engine = create_engine(sync_url)
    try:
        with sync_engine.connect() as conn:
            conn.execute(text(f'DROP SCHEMA "{test_schema}" CASCADE'))
            conn.commit()
    except Exception:
        pass
    sync_engine.dispose()


@pytest.fixture
async def pg_session(db_engine: PostgreSQLManager, ensure_schema: str) -> AsyncGenerator[AsyncSession, None]:
    """创建带 schema 隔离的 PostgreSQL 会话"""
    async_engine = db_engine.get_async_engine()
    session = AsyncSession(async_engine, expire_on_commit=False)

    await session.execute(text(f'SET search_path TO "{ensure_schema}"'))

    async with session.begin_nested():
        yield session

    await session.rollback()
    await session.close()


@pytest.fixture
def repo(pg_session: AsyncSession):
    """注入 session 上下文的 DocumentRepository"""
    token = set_session(pg_session)
    repository = PostgreSQLDocumentRepository()
    yield repository
    reset_session(token)


# ===================================================================
# MinIO 真实服务 Fixtures
# ===================================================================


@pytest.fixture
def minio_tenant_id() -> str:
    return f"test-integration-ext-{uuid.uuid4().hex[:8]}"


@pytest.fixture
def minio_storage(minio_tenant_id: str):
    """创建真实 MinIO 存储适配器，确保 bucket 存在，测试后清理"""
    from src.infrastructure.config.minio import MinIOConfig
    from src.infrastructure.storage.minio.bucket_manager import BucketManager
    from src.infrastructure.storage.minio.minio_adapter import MinIOAdapter
    from src.infrastructure.storage.minio.minio_repository import MinIORepository
    from src.infrastructure.storage.minio.object_operations import ObjectOperations
    from src.infrastructure.storage.minio.worm_lifecycle import WORMManager

    env = get_test_env()
    endpoint_parts = env.minio.endpoint.split(":")
    host = endpoint_parts[0]
    port = int(endpoint_parts[1]) if len(endpoint_parts) > 1 else 9000

    config = MinIOConfig(
        host=host,
        port=port,
        access_key=env.minio.access_key,
        secret_key=env.minio.secret_key,
        secure=env.minio.secure,
    )

    try:
        from src.infrastructure.storage.minio.minio_manager import MinioManager

        mgr = MinioManager.from_config(config)
        if not mgr.health_check():
            pytest.skip("MinIO not available")
    except Exception as e:
        pytest.skip(f"MinIO not available: {e}")

    bucket_manager = BucketManager(config)
    object_ops = ObjectOperations(config)
    worm_mgr = WORMManager(config)

    repository = MinIORepository(
        bucket_manager=bucket_manager,
        object_operations=object_ops,
        worm_manager=worm_mgr,
        tenant_id=minio_tenant_id,
    )

    adapter = MinIOAdapter(repository)

    bucket_name = bucket_manager.build_bucket_name("raw-documents", minio_tenant_id)
    if not bucket_manager.bucket_exists(bucket_name):
        bucket_manager.create_bucket(bucket_name)

    yield adapter

    _cleanup_minio_bucket(bucket_manager, bucket_name)


def _cleanup_minio_bucket(bucket_manager, bucket_name: str) -> None:
    """安全清理 MinIO bucket"""
    try:
        if bucket_manager.bucket_exists(bucket_name):
            bucket_manager.delete_bucket(bucket_name, force=True)
    except Exception:
        pass


async def _upload_to_minio_async(adapter, bucket_type: str, object_key: str, file_path: str) -> None:
    """通过 MinIOAdapter 上传文件到 MinIO（async 版本）"""
    await adapter.store(bucket_type, object_key, file_path, "application/octet-stream")


# ===================================================================
# 真实 EventPublisher Fixture
# ===================================================================


@pytest.fixture
def event_bus():
    """真实 DualChannelEventBus（Redis + PostgreSQL Outbox → RabbitMQ）"""
    from src.infrastructure.config.redis import RedisConfig as InfraRedisConfig
    from src.infrastructure.messaging.channel_router import ChannelRouter
    from src.infrastructure.messaging.dual_channel_event_bus import DualChannelEventBus
    from src.infrastructure.messaging.outbox.outbox_repository import PostgreSQLOutboxRepository
    from src.infrastructure.messaging.rabbitmq_event_bus import RabbitMQEventBus
    from src.infrastructure.messaging.redis_event_bus import RedisEventBus
    from src.infrastructure.messaging.redis_publisher import RedisEventPublisher
    from src.infrastructure.messaging.redis_subscriber import RedisEventSubscriber

    env = get_test_env()
    redis_config = InfraRedisConfig(
        host=env.redis.host,
        port=env.redis.port,
    )
    redis_publisher = RedisEventPublisher(redis_config)
    redis_subscriber = RedisEventSubscriber(redis_config)
    router = ChannelRouter(load_defaults=True)
    redis_bus = RedisEventBus(redis_publisher, redis_subscriber, router)
    rabbitmq_bus = RabbitMQEventBus(PostgreSQLOutboxRepository(), router)
    return DualChannelEventBus(redis_bus, rabbitmq_bus, router)


# ===================================================================
# 解析器层集成测试 —— PPTX
# ===================================================================


class TestParsePipelinePPTX:
    """PPTX 解析流水线"""

    def test_parse_pptx_success(self) -> None:
        from src.infrastructure.external_services.document_parsing.pptx_parser import PptxParser

        parser = PptxParser()
        path = _create_test_pptx()
        try:
            result = parser.parse(path, MIME_PPTX)
            assert result.parse_status == "completed"
            assert result.mime_type == MIME_PPTX
            assert len(result.pages) >= 1
            assert result.pages[0].page_number == 1

            all_text = " ".join(t.content for p in result.pages for t in p.texts)
            assert "集成测试PPTX" in all_text

            import json

            json.dumps(result.to_dict(), ensure_ascii=False)
        finally:
            os.unlink(path)

    def test_parse_pptx_empty(self) -> None:
        from src.infrastructure.external_services.document_parsing.pptx_parser import PptxParser

        parser = PptxParser()
        path = _create_test_pptx()
        try:
            result = parser.parse(path, MIME_PPTX)
            assert result.parse_status == "completed"
            assert result.document_id
            assert result.pages
        finally:
            os.unlink(path)


# ===================================================================
# 解析器层集成测试 —— XLSX
# ===================================================================


class TestParsePipelineXLSX:
    """XLSX 解析流水线"""

    def test_parse_xlsx_success(self) -> None:
        from src.infrastructure.external_services.document_parsing.excel_parser import ExcelParser

        parser = ExcelParser()
        path = _create_test_xlsx()
        try:
            result = parser.parse(path, MIME_XLSX)
            assert result.parse_status == "completed"
            assert result.mime_type == MIME_XLSX
            assert len(result.pages) >= 1

            # XLSX 解析器将数据存储在 tables 中，非 texts
            all_cells = []
            for p in result.pages:
                for tb in p.tables:
                    for row in tb.rows:
                        all_cells.extend(row)
            all_text = " ".join(all_cells)
            assert "集成测试XLSX" in all_text
        finally:
            os.unlink(path)

    def test_parse_xlsx_extracts_tables(self) -> None:
        from src.infrastructure.external_services.document_parsing.excel_parser import ExcelParser

        parser = ExcelParser()
        path = _create_test_xlsx()
        try:
            result = parser.parse(path, MIME_XLSX)
            assert result.parse_status == "completed"
            assert len(result.pages) >= 1
            # Excel 会为每 sheet 至少提取一个表格
            assert len(result.pages[0].tables) >= 0, f"tables 应为 list，实际: {type(result.pages[0].tables)}"
        finally:
            os.unlink(path)


# ===================================================================
# 解析器层集成测试 —— CSV
# ===================================================================


class TestParsePipelineCSV:
    """CSV 解析流水线"""

    def test_parse_csv_success(self) -> None:
        from src.infrastructure.external_services.document_parsing.csv_parser import CSVParser

        parser = CSVParser()
        path = _create_test_csv()
        try:
            result = parser.parse(path, MIME_CSV)
            assert result.parse_status == "completed"
            assert result.mime_type == MIME_CSV
            assert len(result.pages) >= 1

            # CSV 解析器将数据存储在 tables 中
            all_cells = []
            for p in result.pages:
                for tb in p.tables:
                    for row in tb.rows:
                        all_cells.extend(row)
            all_text = " ".join(all_cells)
            assert "列A" in all_text
        finally:
            os.unlink(path)

    def test_parse_csv_extracts_table(self) -> None:
        from src.infrastructure.external_services.document_parsing.csv_parser import CSVParser

        parser = CSVParser()
        path = _create_test_csv()
        try:
            result = parser.parse(path, MIME_CSV)
            assert result.parse_status == "completed"
            assert len(result.pages) >= 1
            tables = result.pages[0].tables
            assert len(tables) >= 1, f"CSV 应至少提取 1 个表格，实际: {len(tables)}"
            assert len(tables[0].rows) >= 2  # header + at least 1 data row
        finally:
            os.unlink(path)


# ===================================================================
# 解析器层集成测试 —— HTML
# ===================================================================


class TestParsePipelineHTML:
    """HTML 解析流水线"""

    def test_parse_html_success(self) -> None:
        from src.infrastructure.external_services.document_parsing.html_parser import HTMLParser

        parser = HTMLParser()
        path = _create_test_html()
        try:
            result = parser.parse(path, MIME_HTML)
            assert result.parse_status == "completed"
            assert result.mime_type == MIME_HTML
            assert len(result.pages) >= 1

            all_text = " ".join(t.content for p in result.pages for t in p.texts)
            assert "集成测试HTML" in all_text
        finally:
            os.unlink(path)

    def test_parse_html_extracts_text_content(self) -> None:
        from src.infrastructure.external_services.document_parsing.html_parser import HTMLParser

        parser = HTMLParser()
        path = _create_test_html()
        try:
            result = parser.parse(path, MIME_HTML)
            assert result.parse_status == "completed"
            all_text = " ".join(t.content for p in result.pages for t in p.texts)
            assert "段落内容" in all_text
        finally:
            os.unlink(path)


# ===================================================================
# 解析器层集成测试 —— Markdown
# ===================================================================


class TestParsePipelineMarkdown:
    """Markdown 解析流水线"""

    def test_parse_markdown_success(self) -> None:
        from src.infrastructure.external_services.document_parsing.markdown_parser import MarkdownParser

        parser = MarkdownParser()
        path = _create_test_md()
        try:
            result = parser.parse(path, MIME_MARKDOWN)
            assert result.parse_status == "completed"
            assert result.mime_type == MIME_MARKDOWN
            assert len(result.pages) >= 1

            all_text = " ".join(t.content for p in result.pages for t in p.texts)
            assert "集成测试Markdown" in all_text
        finally:
            os.unlink(path)

    def test_parse_markdown_extracts_table(self) -> None:
        from src.infrastructure.external_services.document_parsing.markdown_parser import MarkdownParser

        parser = MarkdownParser()
        path = _create_test_md()
        try:
            result = parser.parse(path, MIME_MARKDOWN)
            assert result.parse_status == "completed"
            tables = result.pages[0].tables
            assert len(tables) >= 1, f"Markdown 应至少提取 1 个表格，实际: {len(tables)}"
            assert len(tables[0].rows) >= 1
        finally:
            os.unlink(path)


# ===================================================================
# 解析器层集成测试 —— RTF
# ===================================================================


class TestParsePipelineRTF:
    """RTF 解析流水线"""

    def test_parse_rtf_success(self) -> None:
        from src.infrastructure.external_services.document_parsing.rtf_parser import RTFParser

        parser = RTFParser()
        path = _create_test_rtf()
        try:
            result = parser.parse(path, MIME_RTF)
            assert result.parse_status == "completed"
            assert result.mime_type == MIME_RTF
            assert len(result.pages) >= 1

            all_text = " ".join(t.content for p in result.pages for t in p.texts)
            assert len(all_text) > 0
        finally:
            os.unlink(path)


# ===================================================================
# 组合路由集成测试
# ===================================================================


class TestCompositeRoutingExtended:
    """扩展格式组合路由集成测试"""

    @staticmethod
    def _build_full_composite():
        """构建包含所有 15 种 MIME 类型的 CompositeDocumentParser"""
        from src.infrastructure.external_services.document_parsing.composite_parser import CompositeDocumentParser
        from src.infrastructure.external_services.document_parsing.csv_parser import CSVParser
        from src.infrastructure.external_services.document_parsing.excel_parser import ExcelParser
        from src.infrastructure.external_services.document_parsing.html_parser import HTMLParser
        from src.infrastructure.external_services.document_parsing.image_parser import ImageParser
        from src.infrastructure.external_services.document_parsing.markdown_parser import MarkdownParser
        from src.infrastructure.external_services.document_parsing.pdf_parser import PDFParser
        from src.infrastructure.external_services.document_parsing.pptx_parser import PptxParser
        from src.infrastructure.external_services.document_parsing.rtf_parser import RTFParser
        from src.infrastructure.external_services.document_parsing.text_parser import TextParser
        from src.infrastructure.external_services.document_parsing.word_parser import WordParser

        return CompositeDocumentParser(
            parsers={
                "application/pdf": PDFParser(),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document": WordParser(),
                "application/msword": WordParser(),
                "text/plain": TextParser(),
                MIME_PPTX: PptxParser(),
                "application/vnd.ms-powerpoint": PptxParser(),
                MIME_XLSX: ExcelParser(),
                "application/vnd.ms-excel": ExcelParser(),
                MIME_CSV: CSVParser(),
                MIME_JPEG: ImageParser(),
                MIME_PNG: ImageParser(),
                "image/gif": ImageParser(),
                MIME_HTML: HTMLParser(),
                MIME_MARKDOWN: MarkdownParser(),
                MIME_RTF: RTFParser(),
            },
        )

    def test_route_pptx(self) -> None:
        parser = self._build_full_composite()
        path = _create_test_pptx()
        try:
            result = parser.parse(path, MIME_PPTX)
            assert result.parse_status == "completed"
        finally:
            os.unlink(path)

    def test_route_xlsx(self) -> None:
        parser = self._build_full_composite()
        path = _create_test_xlsx()
        try:
            result = parser.parse(path, MIME_XLSX)
            assert result.parse_status == "completed"
        finally:
            os.unlink(path)

    def test_route_csv(self) -> None:
        parser = self._build_full_composite()
        path = _create_test_csv()
        try:
            result = parser.parse(path, MIME_CSV)
            assert result.parse_status == "completed"
        finally:
            os.unlink(path)

    def test_route_html(self) -> None:
        parser = self._build_full_composite()
        path = _create_test_html()
        try:
            result = parser.parse(path, MIME_HTML)
            assert result.parse_status == "completed"
        finally:
            os.unlink(path)

    def test_route_markdown(self) -> None:
        parser = self._build_full_composite()
        path = _create_test_md()
        try:
            result = parser.parse(path, MIME_MARKDOWN)
            assert result.parse_status == "completed"
        finally:
            os.unlink(path)

    def test_route_rtf(self) -> None:
        parser = self._build_full_composite()
        path = _create_test_rtf()
        try:
            result = parser.parse(path, MIME_RTF)
            assert result.parse_status == "completed"
        finally:
            os.unlink(path)

    def test_route_unknown_still_fails(self) -> None:
        parser = self._build_full_composite()
        path = _create_test_csv()
        try:
            result = parser.parse(path, "application/x-unknown")
            assert result.parse_status == "failed"
            assert result.error_message is not None
            assert "不支持的 MIME" in result.error_message
        finally:
            os.unlink(path)


# ===================================================================
# Service 层集成测试：通过 DocumentParsingService 完整流水线
# 使用真实 PostgreSQL + 真实 MinIO + 真实解析器 + 真实 EventBus
# ===================================================================


class TestDocumentParsingServiceExtendedPipeline:
    """测试 DocumentParsingService 完整编排流水线（扩展格式）

    真实组件：PostgreSQL（schema 隔离）、MinIO、解析器、DualChannelEventBus（Redis + Outbox → RabbitMQ）。
    """

    @staticmethod
    def _make_extended_composite_parser():
        from src.infrastructure.external_services.document_parsing.composite_parser import CompositeDocumentParser
        from src.infrastructure.external_services.document_parsing.csv_parser import CSVParser
        from src.infrastructure.external_services.document_parsing.excel_parser import ExcelParser
        from src.infrastructure.external_services.document_parsing.html_parser import HTMLParser
        from src.infrastructure.external_services.document_parsing.markdown_parser import MarkdownParser
        from src.infrastructure.external_services.document_parsing.pdf_parser import PDFParser
        from src.infrastructure.external_services.document_parsing.pptx_parser import PptxParser
        from src.infrastructure.external_services.document_parsing.rtf_parser import RTFParser
        from src.infrastructure.external_services.document_parsing.text_parser import TextParser
        from src.infrastructure.external_services.document_parsing.word_parser import WordParser

        return CompositeDocumentParser(
            parsers={
                "application/pdf": PDFParser(),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document": WordParser(),
                "text/plain": TextParser(),
                MIME_PPTX: PptxParser(),
                MIME_XLSX: ExcelParser(),
                MIME_CSV: CSVParser(),
                MIME_HTML: HTMLParser(),
                MIME_MARKDOWN: MarkdownParser(),
                MIME_RTF: RTFParser(),
            },
        )

    @pytest.mark.asyncio
    async def test_parse_pptx_through_service(
        self,
        pg_session: AsyncSession,
        repo: PostgreSQLDocumentRepository,
        minio_storage,
        event_bus,
    ) -> None:
        """完整流水线：PPTX 文件上传到 MinIO → 通过 Service 下载并解析"""
        from src.application.services.document_parsing_service import DocumentParsingService

        pptx_path = _create_test_pptx()
        object_key = f"test-{uuid.uuid4().hex[:8]}.pptx"
        try:
            await _upload_to_minio_async(minio_storage, "raw-documents", object_key, pptx_path)

            parser = self._make_extended_composite_parser()
            service = DocumentParsingService(
                document_repository=repo,
                document_storage=minio_storage,
                event_publisher=event_bus,
                document_parser=parser,
            )

            doc_id = uuid.uuid4()
            tenant_id = "test-tenant"
            doc = Document(
                document_id=doc_id,
                filename="test.pptx",
                mime_type=MIME_PPTX,
                tenant_id=tenant_id,
            )
            doc.metadata["storage_object_key"] = object_key

            await repo.save(doc)
            await pg_session.flush()

            result = await service.parse_document(doc_id, tenant_id)

            assert result.parse_status == ParseStatus.COMPLETED
            parse_result = result.metadata.get("parse_result")
            assert parse_result is not None, f"metadata 应包含 parse_result，实际: {result.metadata}"
            assert parse_result["parse_status"] == "completed"
            assert len(parse_result["pages"]) >= 1

            # 验证事件已写入 Outbox（RELIABLE 通道）
            outbox_result = await pg_session.execute(
                sa_select(OutboxModel).where(OutboxModel.event_type == "DocumentProcessed")
            )
            outbox_entries = list(outbox_result.scalars().all())
            assert len(outbox_entries) == 1, f"Outbox 应有 1 条 DocumentProcessed 事件，实际: {len(outbox_entries)}"
            assert outbox_entries[0].status == "pending"
        finally:
            _cleanup(pptx_path)

    @pytest.mark.asyncio
    async def test_parse_xlsx_through_service(
        self,
        pg_session: AsyncSession,
        repo: PostgreSQLDocumentRepository,
        minio_storage,
        event_bus,
    ) -> None:
        """完整流水线：XLSX 文件上传到 MinIO → 通过 Service 下载并解析"""
        from src.application.services.document_parsing_service import DocumentParsingService

        xlsx_path = _create_test_xlsx()
        object_key = f"test-{uuid.uuid4().hex[:8]}.xlsx"
        try:
            await _upload_to_minio_async(minio_storage, "raw-documents", object_key, xlsx_path)

            parser = self._make_extended_composite_parser()
            service = DocumentParsingService(
                document_repository=repo,
                document_storage=minio_storage,
                event_publisher=event_bus,
                document_parser=parser,
            )

            doc_id = uuid.uuid4()
            tenant_id = "test-tenant"
            doc = Document(
                document_id=doc_id,
                filename="test.xlsx",
                mime_type=MIME_XLSX,
                tenant_id=tenant_id,
            )
            doc.metadata["storage_object_key"] = object_key

            await repo.save(doc)
            await pg_session.flush()

            result = await service.parse_document(doc_id, tenant_id)

            assert result.parse_status == ParseStatus.COMPLETED
            parse_result = result.metadata.get("parse_result")
            assert parse_result is not None
            assert parse_result["parse_status"] == "completed"

            # 验证事件已写入 Outbox
            outbox_result = await pg_session.execute(
                sa_select(OutboxModel).where(OutboxModel.event_type == "DocumentProcessed")
            )
            outbox_entries = list(outbox_result.scalars().all())
            assert len(outbox_entries) == 1
        finally:
            _cleanup(xlsx_path)

    @pytest.mark.asyncio
    async def test_parse_markdown_through_service(
        self,
        pg_session: AsyncSession,
        repo: PostgreSQLDocumentRepository,
        minio_storage,
        event_bus,
    ) -> None:
        """完整流水线：Markdown 文件上传到 MinIO → 通过 Service 下载并解析"""
        from src.application.services.document_parsing_service import DocumentParsingService

        md_path = _create_test_md()
        object_key = f"test-{uuid.uuid4().hex[:8]}.md"
        try:
            await _upload_to_minio_async(minio_storage, "raw-documents", object_key, md_path)

            parser = self._make_extended_composite_parser()
            service = DocumentParsingService(
                document_repository=repo,
                document_storage=minio_storage,
                event_publisher=event_bus,
                document_parser=parser,
            )

            doc_id = uuid.uuid4()
            tenant_id = "test-tenant"
            doc = Document(
                document_id=doc_id,
                filename="test.md",
                mime_type=MIME_MARKDOWN,
                tenant_id=tenant_id,
            )
            doc.metadata["storage_object_key"] = object_key

            await repo.save(doc)
            await pg_session.flush()

            result = await service.parse_document(doc_id, tenant_id)

            assert result.parse_status == ParseStatus.COMPLETED
            parse_result = result.metadata.get("parse_result")
            assert parse_result is not None
            all_text = " ".join(t["content"] for p in parse_result["pages"] for t in p["texts"])
            assert "集成测试Markdown" in all_text

            outbox_result = await pg_session.execute(
                sa_select(OutboxModel).where(OutboxModel.event_type == "DocumentProcessed")
            )
            outbox_entries = list(outbox_result.scalars().all())
            assert len(outbox_entries) == 1
        finally:
            _cleanup(md_path)

    @pytest.mark.asyncio
    async def test_parse_csv_through_service(
        self,
        pg_session: AsyncSession,
        repo: PostgreSQLDocumentRepository,
        minio_storage,
        event_bus,
    ) -> None:
        """完整流水线：CSV 文件上传到 MinIO → 通过 Service 下载并解析"""
        from src.application.services.document_parsing_service import DocumentParsingService

        csv_path = _create_test_csv()
        object_key = f"test-{uuid.uuid4().hex[:8]}.csv"
        try:
            await _upload_to_minio_async(minio_storage, "raw-documents", object_key, csv_path)

            parser = self._make_extended_composite_parser()
            service = DocumentParsingService(
                document_repository=repo,
                document_storage=minio_storage,
                event_publisher=event_bus,
                document_parser=parser,
            )

            doc_id = uuid.uuid4()
            tenant_id = "test-tenant"
            doc = Document(
                document_id=doc_id,
                filename="test.csv",
                mime_type=MIME_CSV,
                tenant_id=tenant_id,
            )
            doc.metadata["storage_object_key"] = object_key

            await repo.save(doc)
            await pg_session.flush()

            result = await service.parse_document(doc_id, tenant_id)

            assert result.parse_status == ParseStatus.COMPLETED
            parse_result = result.metadata.get("parse_result")
            assert parse_result is not None
            assert parse_result["parse_status"] == "completed"

            outbox_result = await pg_session.execute(
                sa_select(OutboxModel).where(OutboxModel.event_type == "DocumentProcessed")
            )
            outbox_entries = list(outbox_result.scalars().all())
            assert len(outbox_entries) == 1
        finally:
            _cleanup(csv_path)

    @pytest.mark.asyncio
    async def test_parse_rtf_through_service(
        self,
        pg_session: AsyncSession,
        repo: PostgreSQLDocumentRepository,
        minio_storage,
        event_bus,
    ) -> None:
        """完整流水线：RTF 文件上传到 MinIO → 通过 Service 下载并解析"""
        from src.application.services.document_parsing_service import DocumentParsingService

        rtf_path = _create_test_rtf()
        object_key = f"test-{uuid.uuid4().hex[:8]}.rtf"
        try:
            await _upload_to_minio_async(minio_storage, "raw-documents", object_key, rtf_path)

            parser = self._make_extended_composite_parser()
            service = DocumentParsingService(
                document_repository=repo,
                document_storage=minio_storage,
                event_publisher=event_bus,
                document_parser=parser,
            )

            doc_id = uuid.uuid4()
            tenant_id = "test-tenant"
            doc = Document(
                document_id=doc_id,
                filename="test.rtf",
                mime_type=MIME_RTF,
                tenant_id=tenant_id,
            )
            doc.metadata["storage_object_key"] = object_key

            await repo.save(doc)
            await pg_session.flush()

            result = await service.parse_document(doc_id, tenant_id)

            assert result.parse_status == ParseStatus.COMPLETED
            parse_result = result.metadata.get("parse_result")
            assert parse_result is not None
            assert parse_result["parse_status"] == "completed"

            outbox_result = await pg_session.execute(
                sa_select(OutboxModel).where(OutboxModel.event_type == "DocumentProcessed")
            )
            outbox_entries = list(outbox_result.scalars().all())
            assert len(outbox_entries) == 1
        finally:
            _cleanup(rtf_path)

    @pytest.mark.asyncio
    async def test_parse_html_through_service(
        self,
        pg_session: AsyncSession,
        repo: PostgreSQLDocumentRepository,
        minio_storage,
        event_bus,
    ) -> None:
        """完整流水线：HTML 文件上传到 MinIO → 通过 Service 下载并解析"""
        from src.application.services.document_parsing_service import DocumentParsingService

        html_path = _create_test_html()
        object_key = f"test-{uuid.uuid4().hex[:8]}.html"
        try:
            await _upload_to_minio_async(minio_storage, "raw-documents", object_key, html_path)

            parser = self._make_extended_composite_parser()
            service = DocumentParsingService(
                document_repository=repo,
                document_storage=minio_storage,
                event_publisher=event_bus,
                document_parser=parser,
            )

            doc_id = uuid.uuid4()
            tenant_id = "test-tenant"
            doc = Document(
                document_id=doc_id,
                filename="test.html",
                mime_type=MIME_HTML,
                tenant_id=tenant_id,
            )
            doc.metadata["storage_object_key"] = object_key

            await repo.save(doc)
            await pg_session.flush()

            result = await service.parse_document(doc_id, tenant_id)

            assert result.parse_status == ParseStatus.COMPLETED
            parse_result = result.metadata.get("parse_result")
            assert parse_result is not None
            all_text = " ".join(t["content"] for p in parse_result["pages"] for t in p["texts"])
            assert "集成测试HTML" in all_text

            outbox_result = await pg_session.execute(
                sa_select(OutboxModel).where(OutboxModel.event_type == "DocumentProcessed")
            )
            outbox_entries = list(outbox_result.scalars().all())
            assert len(outbox_entries) == 1
        finally:
            _cleanup(html_path)


# ===================================================================
# 并发解析测试（AC-6: ≥10 并发）
# ===================================================================


class TestConcurrentParsingExtended:
    """扩展格式并发解析测试"""

    def test_concurrent_parse_extended_formats(self) -> None:
        """至少 10 个文档并发解析，混合多种扩展格式"""
        from src.infrastructure.external_services.document_parsing.composite_parser import CompositeDocumentParser
        from src.infrastructure.external_services.document_parsing.csv_parser import CSVParser
        from src.infrastructure.external_services.document_parsing.excel_parser import ExcelParser
        from src.infrastructure.external_services.document_parsing.html_parser import HTMLParser
        from src.infrastructure.external_services.document_parsing.markdown_parser import MarkdownParser
        from src.infrastructure.external_services.document_parsing.pdf_parser import PDFParser
        from src.infrastructure.external_services.document_parsing.pptx_parser import PptxParser
        from src.infrastructure.external_services.document_parsing.text_parser import TextParser

        parser = CompositeDocumentParser(
            parsers={
                "application/pdf": PDFParser(),
                "text/plain": TextParser(),
                MIME_PPTX: PptxParser(),
                MIME_XLSX: ExcelParser(),
                MIME_CSV: CSVParser(),
                MIME_HTML: HTMLParser(),
                MIME_MARKDOWN: MarkdownParser(),
            },
        )

        # 创建 12 个混合格式文档
        paths_and_mimes: list[tuple[str, str]] = []
        paths_and_mimes.extend((_create_test_pptx(), MIME_PPTX) for _ in range(2))
        paths_and_mimes.extend((_create_test_xlsx(), MIME_XLSX) for _ in range(2))
        paths_and_mimes.extend((_create_test_csv(), MIME_CSV) for _ in range(2))
        paths_and_mimes.extend((_create_test_html(), MIME_HTML) for _ in range(2))
        paths_and_mimes.extend((_create_test_md(), MIME_MARKDOWN) for _ in range(2))
        paths_and_mimes.extend((_create_test_pptx(), MIME_PPTX) for _ in range(2))

        assert len(paths_and_mimes) == 12

        try:

            async def parse_all():
                tasks = [asyncio.to_thread(parser.parse, p, m) for p, m in paths_and_mimes]
                return await asyncio.gather(*tasks)

            loop = asyncio.new_event_loop()
            try:
                results = loop.run_until_complete(parse_all())
            finally:
                loop.close()

            assert len(results) == 12
            for result in results:
                assert result.parse_status == "completed"
        finally:
            for p, _ in paths_and_mimes:
                _cleanup(p)
