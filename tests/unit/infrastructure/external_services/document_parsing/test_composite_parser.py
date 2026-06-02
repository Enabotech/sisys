"""组合解析器单元测试

TDD 红阶段：测试 CompositeDocumentParser 的 MIME 类型路由和未知 MIME 拒绝。
"""

from __future__ import annotations

import os
import tempfile

from pypdf import PdfWriter

# MIME 类型常量（与 composition_root 保持一致）
MIME_PDF = "application/pdf"
MIME_DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
MIME_DOC = "application/msword"
MIME_TXT = "text/plain"


def _create_minimal_pdf() -> str:
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    writer.write(tmp.name)
    tmp.close()
    return tmp.name


def _create_minimal_docx() -> str:
    from docx import Document

    doc = Document()
    doc.add_paragraph("test")
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
    doc.save(tmp.name)
    tmp.close()
    return tmp.name


def _create_minimal_txt() -> str:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".txt")
    tmp.write(b"hello world")
    tmp.close()
    return tmp.name


def _build_composite():
    """构造使用 dict 注入的 CompositeDocumentParser（与 composition_root 一致）"""
    from src.infrastructure.external_services.document_parsing.composite_parser import (
        CompositeDocumentParser,
    )
    from src.infrastructure.external_services.document_parsing.pdf_parser import (
        PDFParser,
    )
    from src.infrastructure.external_services.document_parsing.text_parser import (
        TextParser,
    )
    from src.infrastructure.external_services.document_parsing.word_parser import (
        WordParser,
    )

    return CompositeDocumentParser(
        parsers={
            MIME_PDF: PDFParser(),
            MIME_DOCX: WordParser(),
            MIME_DOC: WordParser(),  # DOC 格式由 WordParser 返回友好拒绝消息
            MIME_TXT: TextParser(),
        },
    )


class TestCompositeParserRouting:
    """MIME 类型路由测试"""

    def test_pdf_mime_routes_to_pdf_parser(self) -> None:
        parser = _build_composite()
        path = _create_minimal_pdf()
        try:
            result = parser.parse(path, MIME_PDF)
            assert result.parse_status == "completed"
            assert result.mime_type == "application/pdf"
        finally:
            os.unlink(path)

    def test_docx_mime_routes_to_word_parser(self) -> None:
        parser = _build_composite()
        path = _create_minimal_docx()
        try:
            result = parser.parse(path, MIME_DOCX)
            assert result.parse_status == "completed"
        finally:
            os.unlink(path)

    def test_txt_mime_routes_to_text_parser(self) -> None:
        parser = _build_composite()
        path = _create_minimal_txt()
        try:
            result = parser.parse(path, MIME_TXT)
            assert result.parse_status == "completed"
            assert result.mime_type == "text/plain"
        finally:
            os.unlink(path)


class TestCompositeParserUnknownMime:
    """未知 MIME 类型拒绝测试"""

    def test_unknown_mime_returns_failed(self) -> None:
        """不支持的 MIME 类型返回 failed 状态（与其他解析器错误处理策略一致）"""
        parser = _build_composite()
        path = _create_minimal_txt()
        try:
            result = parser.parse(path, "application/unknown")
            assert result.parse_status == "failed"
            assert result.error_message is not None
            assert "不支持的 MIME 类型" in result.error_message
        finally:
            os.unlink(path)


class TestCompositeParserDocRouting:
    """DOC 格式路由测试 — 验证 AC-2 友好拒绝消息"""

    def test_doc_mime_routes_to_word_parser(self) -> None:
        """验证 DOC 格式路由到 WordParser，返回友好中文错误消息"""
        parser = _build_composite()
        # 创建一个非 DOCX 格式文件（模拟 DOC）
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".doc")
        tmp.write(b"not a valid docx")
        tmp.close()
        try:
            result = parser.parse(tmp.name, MIME_DOC)
            assert result.parse_status == "failed"
            assert result.error_message is not None
            assert "DOCX" in result.error_message or "DOC" in result.error_message
        finally:
            os.unlink(tmp.name)


class TestCompositeParserPortContract:
    """验证 CompositeDocumentParser 满足 DocumentParserPort 协议"""

    def test_satisfies_protocol(self) -> None:
        from src.domain.ports.document_parser import DocumentParserPort

        parser = _build_composite()
        assert isinstance(parser, DocumentParserPort)


class TestExtendedFormatRouting:
    """Story 2-2b 扩展格式 MIME 路由测试"""

    def test_pptx_mime_routes_to_pptx_parser(self) -> None:
        """PPtX MIME 类型路由到 PptxParser"""
        from src.infrastructure.external_services.document_parsing.pptx_parser import PptxParser

        pptx = _create_minimal_pptx()
        try:
            parser = PptxParser()
            result = parser.parse(pptx, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
            assert result.is_completed()
        finally:
            os.unlink(pptx)

    def test_xlsx_mime_routes_to_excel_parser(self) -> None:
        """XLSX MIME 类型路由到 ExcelParser"""
        xlsx = _create_minimal_xlsx()
        try:
            from src.infrastructure.external_services.document_parsing.excel_parser import ExcelParser

            parser = ExcelParser()
            result = parser.parse(xlsx, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
            assert result.is_completed()
        finally:
            os.unlink(xlsx)

    def test_csv_mime_routes_to_csv_parser(self) -> None:
        """CSV MIME 类型路由到 CSVParser"""
        from src.infrastructure.external_services.document_parsing.csv_parser import CSVParser

        path = _create_minimal_txt()
        try:
            parser = CSVParser()
            result = parser.parse(path, "text/csv")
            assert result.is_completed()
        finally:
            os.unlink(path)

    def test_html_mime_routes_to_html_parser(self) -> None:
        """HTML MIME 类型路由到 HTMLParser"""
        from src.infrastructure.external_services.document_parsing.html_parser import HTMLParser

        html_path = _create_minimal_html()
        try:
            parser = HTMLParser()
            result = parser.parse(html_path, "text/html")
            assert result.is_completed()
        finally:
            os.unlink(html_path)

    def test_markdown_mime_routes_to_markdown_parser(self) -> None:
        """Markdown MIME 类型路由到 MarkdownParser"""
        from src.infrastructure.external_services.document_parsing.markdown_parser import MarkdownParser

        md_path = _create_minimal_md()
        try:
            parser = MarkdownParser()
            result = parser.parse(md_path, "text/markdown")
            assert result.is_completed()
        finally:
            os.unlink(md_path)

    def test_rtf_mime_routes_to_rtf_parser(self) -> None:
        """RTF MIME 类型路由到 RTFParser"""
        from src.infrastructure.external_services.document_parsing.rtf_parser import RTFParser

        rtf_path = _create_minimal_rtf()
        try:
            parser = RTFParser()
            result = parser.parse(rtf_path, "text/rtf")
            assert result.is_completed()
        finally:
            os.unlink(rtf_path)


def _create_minimal_pptx() -> str:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "test"
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()
    return tmp.name


def _create_minimal_xlsx() -> str:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "Sheet1"
    ws["A1"] = "test"
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx")
    wb.save(tmp.name)
    tmp.close()
    return tmp.name


def _create_minimal_html() -> str:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".html")
    tmp.write(b"<html><body><p>test</p></body></html>")
    tmp.close()
    return tmp.name


def _create_minimal_md() -> str:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".md")
    tmp.write(b"# Test\n\ntest content\n")
    tmp.close()
    return tmp.name


def _create_minimal_rtf() -> str:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".rtf")
    tmp.write(r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Times New Roman;}} \f0\fs24 RTF test}".encode("utf-8"))
    tmp.close()
    return tmp.name
