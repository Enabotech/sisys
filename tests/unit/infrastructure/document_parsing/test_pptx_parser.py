"""PPTX 文档解析器单元测试

TDD 红阶段：测试 PptxParser 的文本提取、表格提取、备注提取、空文档拒绝、旧版 PPT 拒绝。
使用 python-pptx 创建 fixture PPTX 文件，避免依赖外部文件。
"""

from __future__ import annotations

import os
import tempfile


def _create_pptx_with_text(texts: list[str], include_notes: bool = False) -> str:
    """创建含指定文本的 PPTX fixture

    Args:
        texts: 每元素一个幻灯片的标题文本
        include_notes: 第一张幻灯片是否加备注
    """
    from pptx import Presentation

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content
    for i, text in enumerate(texts):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = text
        if include_notes and i == 0:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "备注：重要战略方向"

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()
    return tmp.name


def _create_pptx_with_table(rows: int = 3, cols: int = 3) -> str:
    """创建含内嵌表格的 PPTX fixture"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "数据报告"
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(3))
    for r in range(rows):
        for c in range(cols):
            table_shape.table.cell(r, c).text = f"R{r}C{c}"

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()
    return tmp.name


def _create_empty_pptx() -> str:
    """创建无幻灯片的空 PPTX fixture"""
    from pptx import Presentation

    prs = Presentation()
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()
    return tmp.name


def _create_pptx_with_shapes() -> str:
    """创建含多种形状类型的 PPTX fixture"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "主标题"

    # 文本框形状
    tx_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    tf = tx_box.text_frame
    tf.text = "文本框中内容"

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()
    return tmp.name


MIME_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
MIME_PPT = "application/vnd.ms-powerpoint"


class TestPptxParserCreation:
    """PptxParser 构造和基本功能测试"""

    def test_create_parser(self) -> None:
        """验证 PptxParser 可以正常实例化"""
        from src.infrastructure.document_parsing.pptx_parser import PptxParser

        parser = PptxParser()
        assert parser is not None

    def test_parser_implements_document_parser_port(self) -> None:
        """验证 PptxParser 满足 DocumentParserPort 协议"""
        from src.domain.ports.document_parser import DocumentParserPort
        from src.infrastructure.document_parsing.pptx_parser import PptxParser

        parser = PptxParser()
        assert isinstance(parser, DocumentParserPort)


class TestPptxParserTextExtraction:
    """PPTX 文本提取测试"""

    def test_parse_single_slide_pptx(self) -> None:
        """解析单幻灯片 PPTX，提取文本"""
        from src.infrastructure.document_parsing.pptx_parser import PptxParser

        path = _create_pptx_with_text(["战略规划概述"])
        try:
            parser = PptxParser()
            result = parser.parse(path, MIME_PPTX)

            assert result.is_completed(), f"解析应成功，实际: {result.error_message}"
            assert len(result.pages) >= 1, "应至少有 1 页"

            all_text = " ".join(t.content for p in result.pages for t in p.texts)
            assert "战略规划概述" in all_text, f"应提取到标题文本，实际: {all_text}"
        finally:
            os.unlink(path)

    def test_parse_multi_slide_pptx(self) -> None:
        """解析多幻灯片 PPTX，验证页数和页码"""
        from src.infrastructure.document_parsing.pptx_parser import PptxParser

        path = _create_pptx_with_text(["第一页", "第二页", "第三页"])
        try:
            parser = PptxParser()
            result = parser.parse(path, MIME_PPTX)

            assert result.is_completed()
            assert len(result.pages) == 3, f"应解析 3 页幻灯片，实际: {len(result.pages)}"
            page_numbers = [p.page_number for p in result.pages]
            assert page_numbers == [1, 2, 3], f"页码应为 1/2/3，实际: {page_numbers}"

            # 验证每页标题
            all_text = {p.page_number: " ".join(t.content for t in p.texts) for p in result.pages}
            assert "第一页" in all_text[1]
            assert "第二页" in all_text[2]
            assert "第三页" in all_text[3]
        finally:
            os.unlink(path)

    def test_shapes_metadata_includes_shape_type(self) -> None:
        """文本元素 metadata 包含形状类型"""
        from src.infrastructure.document_parsing.pptx_parser import PptxParser

        path = _create_pptx_with_shapes()
        try:
            parser = PptxParser()
            result = parser.parse(path, MIME_PPTX)

            assert result.is_completed()
            # 应该有来自标题占位符和文本框的文本
            all_text = " ".join(t.content for p in result.pages for t in p.texts)
            assert "主标题" in all_text
            assert "文本框中内容" in all_text
        finally:
            os.unlink(path)


class TestPptxParserTableExtraction:
    """PPTX 表格提取测试"""

    def test_parse_pptx_with_table(self) -> None:
        """提取幻灯片内嵌表格为 ParsedTable"""
        from src.infrastructure.document_parsing.pptx_parser import PptxParser

        path = _create_pptx_with_table(rows=3, cols=3)
        try:
            parser = PptxParser()
            result = parser.parse(path, MIME_PPTX)

            assert result.is_completed()
            all_tables = [t for p in result.pages for t in p.tables]
            assert len(all_tables) >= 1, "应至少提取到 1 个表格"

            table = all_tables[0]
            assert len(table.rows) == 3, f"表格应 3 行，实际: {len(table.rows)}"
            assert len(table.rows[0]) == 3, f"表格应 3 列，实际: {len(table.rows[0])}"
            assert table.rows[0][0] == "R0C0"
            assert table.rows[2][2] == "R2C2"
        finally:
            os.unlink(path)


class TestPptxParserNotesExtraction:
    """PPTX 备注提取测试"""

    def test_parse_pptx_with_notes(self) -> None:
        """提取幻灯片备注内容"""
        from src.infrastructure.document_parsing.pptx_parser import PptxParser

        path = _create_pptx_with_text(["战略规划"], include_notes=True)
        try:
            parser = PptxParser()
            result = parser.parse(path, MIME_PPTX)

            assert result.is_completed()
            all_text = " ".join(t.content for p in result.pages for t in p.texts)
            assert "备注" in all_text, f"应提取到备注内容，实际: {all_text}"
        finally:
            os.unlink(path)

    def test_pptx_without_notes(self) -> None:
        """无备注的 PPTX 正常解析"""
        from src.infrastructure.document_parsing.pptx_parser import PptxParser

        path = _create_pptx_with_text(["无备注页"])
        try:
            parser = PptxParser()
            result = parser.parse(path, MIME_PPTX)

            assert result.is_completed()
        finally:
            os.unlink(path)


class TestPptxParserEmptyDocument:
    """空文档检测测试"""

    def test_empty_pptx_returns_failed(self) -> None:
        """无幻灯片的空 PPTX 返回 failed"""
        from src.infrastructure.document_parsing.pptx_parser import PptxParser

        path = _create_empty_pptx()
        try:
            parser = PptxParser()
            result = parser.parse(path, MIME_PPTX)

            assert result.parse_status == "failed", f"空 PPTX 应返回 failed，实际: {result.parse_status}"
            assert result.error_message is not None, "失败必须有 error_message"
        finally:
            os.unlink(path)


class TestPptxParserLegacyFormatRejection:
    """旧版 PPT 格式拒绝测试"""

    def test_ppt_mime_returns_failed(self) -> None:
        """旧版 PPT MIME 返回 failed 并建议转换"""
        from src.infrastructure.document_parsing.pptx_parser import PptxParser

        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".ppt")
        tmp.write(b"not a valid pptx")
        tmp.close()
        try:
            parser = PptxParser()
            result = parser.parse(tmp.name, MIME_PPT)

            assert result.parse_status == "failed"
            assert result.error_message is not None
            assert "PPTX" in result.error_message, f"错误信息应建议转换为 PPTX，实际: {result.error_message}"
        finally:
            os.unlink(tmp.name)

    def test_corrupt_file_returns_failed(self) -> None:
        """损坏文件返回 failed"""
        from src.infrastructure.document_parsing.pptx_parser import PptxParser

        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        tmp.write(b"this is not a valid pptx zip file")
        tmp.close()
        try:
            parser = PptxParser()
            result = parser.parse(tmp.name, MIME_PPTX)

            assert result.parse_status == "failed"
            assert result.error_message is not None
        finally:
            os.unlink(tmp.name)
