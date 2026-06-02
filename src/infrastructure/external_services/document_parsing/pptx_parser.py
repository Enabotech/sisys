"""PPTX 文档解析器

使用 python-pptx 提取幻灯片文本、表格和备注内容的解析器实现。
"""

from __future__ import annotations

import logging
import os
import uuid
from datetime import UTC, datetime

from src.domain.ports.document_parser import DocumentParserPort
from src.domain.value_objects.parsed_document import (
    ParsedDocument,
    ParsedElement,
    ParsedPage,
    ParsedTable,
)
from src.infrastructure.external_services.document_parsing._limits import MAX_PPTX_BYTES

logger = logging.getLogger(__name__)


class PptxParser(DocumentParserPort):
    """PPTX 文档解析器

    使用 python-pptx.Presentation 提取幻灯片内容，支持：
    - 幻灯片文本提取（含形状类型元数据）
    - 内嵌表格提取
    - 备注内容提取
    - 旧版 PPT 格式拒绝（建议转换为 PPTX）
    - 文件大小上限保护
    """

    _PPT_MIME = "application/vnd.ms-powerpoint"

    def parse(self, file_path: str, mime_type: str) -> ParsedDocument:
        """解析 PPTX 文件

        Args:
            file_path: 本地 PPTX 文件路径
            mime_type: MIME 类型（用于 PPT 格式拒绝和路由决策）

        Returns:
            结构化解析结果
        """
        doc_id = str(uuid.uuid4())
        timestamp = datetime.now(UTC).isoformat()

        if mime_type == self._PPT_MIME:
            return ParsedDocument(
                document_id=doc_id,
                mime_type=mime_type,
                parse_status="failed",
                error_message="不支持旧版 PPT 格式，请转换为 PPTX",
                parse_timestamp=timestamp,
            )

        try:
            file_size = os.path.getsize(file_path)
        except OSError:
            logger.exception("PPTX 文件大小检查失败")
            return ParsedDocument(
                document_id=doc_id,
                mime_type=mime_type,
                parse_status="failed",
                error_message="无法访问文件，请检查文件路径或权限",
                parse_timestamp=timestamp,
            )

        if file_size == 0:
            return ParsedDocument(
                document_id=doc_id,
                mime_type=mime_type,
                parse_status="failed",
                error_message="PPTX 文档为空",
                parse_timestamp=timestamp,
            )

        if file_size > MAX_PPTX_BYTES:
            return ParsedDocument(
                document_id=doc_id,
                mime_type=mime_type,
                parse_status="failed",
                error_message=f"PPTX 文件大小 {file_size // (1024 * 1024)}MB 超过 {MAX_PPTX_BYTES // (1024 * 1024)}MB 限制",
                parse_timestamp=timestamp,
            )

        try:
            from pptx import Presentation

            with open(file_path, "rb") as f:
                prs = Presentation(f)

            slides = list(prs.slides)
            if len(slides) == 0:
                return ParsedDocument(
                    document_id=doc_id,
                    mime_type=mime_type,
                    parse_status="failed",
                    error_message="PPTX 文档为空，无幻灯片",
                    parse_timestamp=timestamp,
                )

            pages: list[ParsedPage] = []
            for slide_idx, slide in enumerate(slides):
                page_number = slide_idx + 1  # 1-indexed
                texts: list[ParsedElement] = []
                tables: list[ParsedTable] = []

                for shape in slide.shapes:
                    shape_type = shape.shape_type
                    shape_type_name = str(shape_type) if shape_type is not None else "UNKNOWN"

                    # 文本提取
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                texts.append(
                                    ParsedElement(
                                        content=text,
                                        metadata={
                                            "shape_type": shape_type_name,
                                            "shape_name": shape.name,
                                        },
                                    )
                                )

                    # 表格提取
                    if shape.has_table:
                        table = shape.table
                        rows_data: list[list[str]] = []
                        for row in table.rows:
                            row_data = [cell.text for cell in row.cells]
                            rows_data.append(row_data)
                        tables.append(ParsedTable(rows=rows_data))

                # 备注提取
                try:
                    notes_slide = slide.notes_slide
                    for paragraph in notes_slide.notes_text_frame.paragraphs:
                        note_text = paragraph.text.strip()
                        if note_text:
                            texts.append(
                                ParsedElement(
                                    content=note_text,
                                    metadata={"shape_type": "NOTES", "shape_name": "slide_notes"},
                                )
                            )
                except (ValueError, AttributeError):
                    # 幻灯片无备注时 notes_slide 不存在或引发异常
                    pass

                if texts or tables:
                    pages.append(
                        ParsedPage(
                            page_number=page_number,
                            texts=texts,
                            tables=tables,
                        )
                    )

            if not pages:
                return ParsedDocument(
                    document_id=doc_id,
                    mime_type=mime_type,
                    parse_status="failed",
                    error_message="PPTX 文档无有效文本或表格内容",
                    parse_timestamp=timestamp,
                )

            return ParsedDocument(
                document_id=doc_id,
                mime_type=mime_type,
                pages=pages,
                parse_status="completed",
                parse_timestamp=timestamp,
            )

        except Exception:
            logger.exception("PPTX 解析失败")
            return ParsedDocument(
                document_id=doc_id,
                mime_type=mime_type,
                parse_status="failed",
                error_message="PPTX 文档解析失败，文件可能已损坏或格式不正确",
                parse_timestamp=timestamp,
            )
