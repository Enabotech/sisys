"""Office 文档格式处理器单元测试

验证 OfficeFormatHandler DOCX/PPTX/XLSX 内容标题提取

"""

from __future__ import annotations

import tempfile
from pathlib import Path

from plugins.crawler.core.format.handlers.office_handler import OfficeFormatHandler
from plugins.crawler.core.value_objects import FileMetadata


class TestOfficeFormatHandler:
    """OfficeFormatHandler 测试"""

    def setup_method(self) -> None:
        self.handler = OfficeFormatHandler()

    def test_can_handle_docx(self) -> None:
        """应识别 .docx 扩展名"""
        assert self.handler.can_handle("test.docx", "")

    def test_can_handle_pptx(self) -> None:
        """应识别 .pptx 扩展名"""
        assert self.handler.can_handle("test.pptx", "")

    def test_can_handle_xlsx(self) -> None:
        """应识别 .xlsx 扩展名"""
        assert self.handler.can_handle("test.xlsx", "")

    def test_cannot_handle_pdf(self) -> None:
        """不应处理 PDF"""
        assert not self.handler.can_handle("test.pdf", "application/pdf")

    def test_unsupported_ext_returns_empty(self) -> None:
        """不支持的扩展名应返回空元数据"""
        meta = self.handler.extract_metadata("test.doc")
        assert meta == FileMetadata()


class TestDocxContentTitle:
    """DOCX 内容标题提取测试"""

    def setup_method(self) -> None:
        self.handler = OfficeFormatHandler()

    def _create_docx(
        self,
        heading: str = "",
        first_paragraph: str = "",
        title: str = "",
        author: str = "",
    ) -> str:
        """创建测试 DOCX 文件

        Args:
            heading: 标题样式段落文本
            first_paragraph: 首个普通段落文本
            title: 文档元数据标题
            author: 文档作者

        Returns:
            临时文件路径
        """
        from docx import Document

        doc = Document()

        if title or author:
            doc.core_properties.title = title
            doc.core_properties.author = author

        if heading:
            doc.add_heading(heading, level=1)

        if first_paragraph:
            doc.add_paragraph(first_paragraph)

        tmp = tempfile.NamedTemporaryFile(suffix=".docx", delete=False)
        doc.save(tmp.name)
        tmp.close()
        return tmp.name

    def test_heading_style_extracted(self) -> None:
        """应从 Heading 样式提取内容标题"""
        path = self._create_docx(heading="战略规划报告")
        try:
            meta = self.handler.extract_metadata(path)
            assert "战略规划报告" in meta.content_title
        finally:
            Path(path).unlink(missing_ok=True)

    def test_fallback_to_first_paragraph(self) -> None:
        """无 Heading 样式时应降级到首个非空段落"""
        path = self._create_docx(first_paragraph="这是一份重要的报告文档")
        try:
            meta = self.handler.extract_metadata(path)
            assert "报告文档" in meta.content_title
        finally:
            Path(path).unlink(missing_ok=True)

    def test_metadata_title_separate(self) -> None:
        """元数据标题应独立于内容标题"""
        path = self._create_docx(title="文档属性标题", heading="正文标题")
        try:
            meta = self.handler.extract_metadata(path)
            assert meta.title == "文档属性标题"
            assert "正文标题" in meta.content_title
        finally:
            Path(path).unlink(missing_ok=True)

    def test_author_extracted(self) -> None:
        """应提取作者"""
        path = self._create_docx(author="李四", heading="测试")
        try:
            meta = self.handler.extract_metadata(path)
            assert meta.author == "李四"
        finally:
            Path(path).unlink(missing_ok=True)

    def test_empty_doc_returns_empty(self) -> None:
        """空文档应返回空元数据"""
        path = self._create_docx()
        try:
            meta = self.handler.extract_metadata(path)
            assert meta.content_title == ""
        finally:
            Path(path).unlink(missing_ok=True)


class TestPptxContentTitle:
    """PPTX 内容标题提取测试"""

    def setup_method(self) -> None:
        self.handler = OfficeFormatHandler()

    def _create_pptx(
        self,
        title_text: str = "",
        metadata_title: str = "",
    ) -> str:
        """创建测试 PPTX 文件

        Args:
            title_text: 标题占位符文本
            metadata_title: 文档元数据标题

        Returns:
            临时文件路径
        """
        from pptx import Presentation

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
        slide = prs.slides.add_slide(slide_layout)

        if title_text:
            for shape in slide.shapes:
                if shape.has_text_frame and hasattr(shape, "placeholder_format"):
                    shape.text_frame.text = title_text
                    break

        if metadata_title:
            prs.core_properties.title = metadata_title

        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        prs.save(tmp.name)
        tmp.close()
        return tmp.name

    def test_title_placeholder_extracted(self) -> None:
        """应从标题占位符提取内容标题"""
        path = self._create_pptx(title_text="产品发布计划")
        try:
            meta = self.handler.extract_metadata(path)
            assert "产品发布计划" in meta.content_title
        finally:
            Path(path).unlink(missing_ok=True)

    def test_empty_pptx_returns_empty(self) -> None:
        """空演示文稿应返回空内容标题"""
        path = self._create_pptx()
        try:
            meta = self.handler.extract_metadata(path)
            assert meta.content_title == ""
        finally:
            Path(path).unlink(missing_ok=True)

    def test_metadata_title_separate(self) -> None:
        """元数据标题应独立于内容标题"""
        path = self._create_pptx(title_text="幻灯片标题", metadata_title="属性标题")
        try:
            meta = self.handler.extract_metadata(path)
            assert meta.title == "属性标题"
        finally:
            Path(path).unlink(missing_ok=True)


class TestXlsxContentTitle:
    """XLSX 内容标题提取测试"""

    def setup_method(self) -> None:
        self.handler = OfficeFormatHandler()

    def _create_xlsx(
        self,
        sheet_name: str = "Sheet",
        a1_value: str = "",
        metadata_title: str = "",
    ) -> str:
        """创建测试 XLSX 文件

        Args:
            sheet_name: 工作表名
            a1_value: A1 单元格值
            metadata_title: 文档元数据标题

        Returns:
            临时文件路径
        """
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        assert ws is not None
        ws.title = sheet_name

        if a1_value:
            ws["A1"] = a1_value

        if metadata_title:
            wb.properties.title = metadata_title

        tmp = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
        wb.save(tmp.name)
        tmp.close()
        return tmp.name

    def test_meaningful_sheet_name_extracted(self) -> None:
        """应从有意义的工作表名提取标题"""
        path = self._create_xlsx(sheet_name="2024销售数据")
        try:
            meta = self.handler.extract_metadata(path)
            assert "2024销售数据" in meta.content_title
        finally:
            Path(path).unlink(missing_ok=True)

    def test_sheet1_ignored_fallback_to_a1(self) -> None:
        """Sheet1 应被忽略，降级到 A1 值"""
        path = self._create_xlsx(sheet_name="Sheet1", a1_value="员工名单")
        try:
            meta = self.handler.extract_metadata(path)
            assert "员工名单" in meta.content_title
        finally:
            Path(path).unlink(missing_ok=True)

    def test_chinese_sheet1_ignored(self) -> None:
        """中文默认工作表名应被忽略"""
        path = self._create_xlsx(sheet_name="工作表1", a1_value="预算表")
        try:
            meta = self.handler.extract_metadata(path)
            assert "预算表" in meta.content_title
        finally:
            Path(path).unlink(missing_ok=True)

    def test_empty_xlsx_returns_empty(self) -> None:
        """空工作簿应返回空内容标题"""
        path = self._create_xlsx(sheet_name="Sheet1")
        try:
            meta = self.handler.extract_metadata(path)
            assert meta.content_title == ""
        finally:
            Path(path).unlink(missing_ok=True)

    def test_metadata_title_separate(self) -> None:
        """元数据标题应独立于内容标题"""
        path = self._create_xlsx(sheet_name="数据汇总", metadata_title="属性标题")
        try:
            meta = self.handler.extract_metadata(path)
            assert meta.title == "属性标题"
            assert "数据汇总" in meta.content_title
        finally:
            Path(path).unlink(missing_ok=True)
