"""Office 文档格式处理器模块

使用 python-docx / python-pptx / openpyxl 提取 DOCX/PPTX/XLSX 元数据

Author:
    agimtech <agimtech@126.com>

Copyright:
    Copyright (c) 2025-2026 AGIMTECH. All rights reserved.

"""

from __future__ import annotations

import os

from plugins.crawler.core.value_objects import FileMetadata


class OfficeFormatHandler:
    """Office 文档格式处理器 — 支持 DOC/DOCX/PPT/PPTX/XLS/XLSX"""

    EXTENSIONS: tuple[str, ...] = ("doc", "docx", "ppt", "pptx", "xls", "xlsx")
    MIME_TYPES: tuple[str, ...] = (
        "application/msword",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )
    _MEANINGLESS_SHEET_NAMES: frozenset[str] = frozenset(
        {
            "Sheet1",
            "Sheet2",
            "Sheet3",
            "Sheet",
            "工作表1",
            "工作表2",
            "工作表3",
            "工作表",
        }
    )

    @property
    def supported_extensions(self) -> tuple[str, ...]:
        return self.EXTENSIONS

    @property
    def supported_mime_types(self) -> tuple[str, ...]:
        return self.MIME_TYPES

    def can_handle(self, file_path: str, mime_type: str) -> bool:
        ext = os.path.splitext(file_path)[1].lower().lstrip(".")
        return ext in self.EXTENSIONS or mime_type.lower() in self.MIME_TYPES

    def extract_metadata(self, file_path: str) -> FileMetadata:
        """从 Office 文档提取元数据

        Args:
            file_path: 文档路径

        Returns:
            文件元数据
        """
        ext = os.path.splitext(file_path)[1].lower()

        if ext in (".docx",):
            return self._extract_docx(file_path)
        elif ext in (".pptx",):
            return self._extract_pptx(file_path)
        elif ext in (".xlsx",):
            return self._extract_xlsx(file_path)
        else:
            return FileMetadata()

    def _extract_docx(self, file_path: str) -> FileMetadata:
        """提取 DOCX 元数据"""
        try:
            from docx import Document

            doc = Document(file_path)
            props = doc.core_properties

            content_title = self._extract_docx_content_title(doc)

            return FileMetadata(
                title=props.title or "",
                content_title=content_title,
                author=props.author or "",
                subject=props.subject or "",
                created=str(props.created) if props.created else "",
            )
        except Exception:
            return FileMetadata()

    def _extract_docx_content_title(self, doc) -> str:
        """从 DOCX 文档提取内容标题

        优先级：Heading 1/2 样式 > 首个非空段落

        Args:
            doc: Document 实例

        Returns:
            内容推导标题
        """
        try:
            for para in doc.paragraphs:
                if para.style.name.startswith(("Heading 1", "Heading 2", "标题 1", "标题 2")):
                    text: str = para.text.strip()
                    if text:
                        return text[:100]

            for para in doc.paragraphs:
                text = para.text.strip()
                if text and len(text) > 2:
                    return text[:100]

            return ""
        except Exception:
            return ""

    def _extract_pptx(self, file_path: str) -> FileMetadata:
        """提取 PPTX 元数据"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            props = prs.core_properties

            content_title = self._extract_pptx_content_title(prs)

            return FileMetadata(
                title=props.title or "",
                content_title=content_title,
                author=props.author or "",
                subject=props.subject or "",
                created=str(props.created) if props.created else "",
            )
        except Exception:
            return FileMetadata()

    def _extract_pptx_content_title(self, prs) -> str:
        """从 PPTX 首张幻灯片提取标题

        Args:
            prs: Presentation 实例

        Returns:
            内容推导标题
        """
        try:
            if not prs.slides:
                return ""

            first_slide = prs.slides[0]

            for shape in first_slide.shapes:
                if shape.has_text_frame and hasattr(shape, "placeholder_format"):
                    ph_idx = shape.placeholder_format.idx
                    if ph_idx is not None and int(ph_idx) <= 2:
                        return str(shape.text or "").strip()[:100]

            for shape in first_slide.shapes:
                if shape.has_text_frame:
                    text = str(shape.text or "").strip()
                    if text and len(text) > 2:
                        return text[:100]

            return ""
        except Exception:
            return ""

    def _extract_xlsx(self, file_path: str) -> FileMetadata:
        """提取 XLSX 元数据"""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(file_path, read_only=True)
            props = wb.properties

            content_title = self._extract_xlsx_content_title(wb)

            result = FileMetadata(
                title=props.title or "",
                content_title=content_title,
                author=props.creator or "",
                subject=props.subject or "",
                created=str(props.created) if props.created else "",
            )
            wb.close()
            return result
        except Exception:
            return FileMetadata()

    def _extract_xlsx_content_title(self, wb) -> str:
        """从 XLSX 工作簿提取内容标题

        优先级：有意义的工作表名 > A1 单元格值

        Args:
            wb: Workbook 实例

        Returns:
            内容推导标题
        """
        try:
            if not wb.sheetnames:
                return ""

            first_sheet_name = wb.sheetnames[0]

            if first_sheet_name not in self._MEANINGLESS_SHEET_NAMES:
                return str(first_sheet_name)[:100]

            ws = wb[first_sheet_name]
            a1_value = ws["A1"].value
            if a1_value and isinstance(a1_value, str):
                return str(a1_value).strip()[:100]

            return ""
        except Exception:
            return ""
